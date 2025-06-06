# file_extractor.py - ファイル内容抽出機能
import os
import logging
import subprocess
import re
import traceback
from datetime import datetime
import io

logger = logging.getLogger(__name__)

class FileExtractor:
    """ファイル内容抽出クラス - 様々な形式のファイルからテキストを抽出する"""
    
    def __init__(self):
        """FileExtractorの初期化"""
        # 外部ライブラリのインポート状態を追跡
        self.imports = {
            'pdf': False,
            'docx': False,
            'xlsx': False,
            'pptx': False,
        }
        
        # 各種ライブラリの依存関係確認
        self._check_imports()

    def _check_imports(self):
        """利用可能なライブラリをチェック"""
        # PyPDF (PDF抽出用)
        try:
            import PyPDF2
            self.imports['pdf'] = True
            logger.info("PyPDF2が利用可能です - PDFの抽出に使用します")
        except ImportError:
            logger.warning("PyPDF2がインストールされていません。'pip install PyPDF2'でインストールしてください")
            
        # python-docx (Word抽出用)
        try:
            import docx
            self.imports['docx'] = True
            logger.info("python-docxが利用可能です - Word文書の抽出に使用します")
        except ImportError:
            logger.warning("python-docxがインストールされていません。'pip install python-docx'でインストールしてください")
            
        # openpyxl (Excel抽出用)
        try:
            import openpyxl
            self.imports['xlsx'] = True
            logger.info("openpyxlが利用可能です - Excelの抽出に使用します")
        except ImportError:
            logger.warning("openpyxlがインストールされていません。'pip install openpyxl'でインストールしてください")
            
        # python-pptx (PowerPoint抽出用)
        try:
            import pptx
            self.imports['pptx'] = True
            logger.info("python-pptxが利用可能です - PowerPointの抽出に使用します")
        except ImportError:
            logger.warning("python-pptxがインストールされていません。'pip install python-pptx'でインストールしてください")

    def extract_file_content(self, file_path):
        """
        ファイルの内容を抽出する

        Args:
            file_path: 抽出するファイルパス

        Returns:
            ファイルの内容（文字列）
        """
        try:
            # ファイルの存在確認
            if not os.path.exists(file_path):
                return f"ファイル '{os.path.basename(file_path)}' が見つかりません。削除または移動された可能性があります。"

            # ファイルサイズ確認 (100MB以上は処理しない)
            file_size = os.path.getsize(file_path)
            if file_size > 100 * 1024 * 1024:  # 100MB
                return f"ファイル '{os.path.basename(file_path)}' は{file_size / (1024 * 1024):.1f}MBと大きすぎるため、処理できません。"
            
            # アクセス権確認
            if not os.access(file_path, os.R_OK):
                return f"ファイル '{os.path.basename(file_path)}' へのアクセス権限がありません。"

            # ファイルの拡張子を取得
            _, ext = os.path.splitext(file_path.lower())

            # ファイルタイプに応じた抽出処理
            if ext == '.pdf':
                return self._extract_pdf(file_path)
            elif ext == '.docx':
                return self._extract_docx(file_path)
            elif ext == '.xlsx':
                return self._extract_xlsx(file_path)
            elif ext == '.pptx':
                return self._extract_pptx(file_path)
            elif ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm', '.log', '.py', '.js', '.css']:
                return self._extract_text(file_path)
            else:
                # 未対応のファイル形式
                file_info = self._get_file_info(file_path)
                return f"未対応のファイル形式 ({ext}):\n{file_info}"

        except PermissionError:
            logger.error(f"ファイル '{file_path}' へのアクセス権限がありません")
            return f"ファイル '{os.path.basename(file_path)}' へのアクセス権限がありません。システム管理者に確認してください。"
        except Exception as e:
            logger.error(f"ファイル '{file_path}' の抽出中にエラーが発生しました: {str(e)}")
            logger.error(traceback.format_exc())
            return f"ファイル抽出エラー: {str(e)}"

    def _extract_text(self, file_path):
        """テキストファイルの内容を抽出"""
        try:
            # 複数のエンコーディングを試行
            encodings = ['utf-8', 'shift-jis', 'cp932', 'euc-jp', 'iso-2022-jp']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                        content = f.read()
                        
                        # ファイル情報のヘッダーを追加
                        file_info = self._get_file_info(file_path)
                        return f"{file_info}\n\n{content}"
                except UnicodeDecodeError:
                    continue
            
            # すべてのエンコーディングで失敗した場合
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
                file_info = self._get_file_info(file_path)
                return f"{file_info}\n\n{content} (エンコーディングの問題があるため、一部文字化けしている可能性があります)"
                
        except Exception as e:
            logger.error(f"テキストファイル '{file_path}' の読み込み中にエラー: {str(e)}")
            return f"テキストファイル読み込みエラー: {str(e)}"

    def _extract_pdf(self, file_path):
        """PDFファイルからテキストを抽出"""
        file_info = self._get_file_info(file_path)
        
        # PyPDF2が利用可能な場合
        if self.imports['pdf']:
            try:
                import PyPDF2
                
                text_content = []
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    
                    # PDF基本情報
                    info = reader.metadata
                    if info:
                        text_content.append(f"タイトル: {info.get('/Title', '不明')}")
                        text_content.append(f"作成者: {info.get('/Author', '不明')}")
                        text_content.append(f"作成日: {info.get('/CreationDate', '不明')}")
                    
                    # ページ数
                    num_pages = len(reader.pages)
                    text_content.append(f"ページ数: {num_pages}")
                    text_content.append("----------------------------------------")
                    
                    # 各ページのテキストを抽出
                    for i, page in enumerate(reader.pages):
                        if i < 10:  # 最初の10ページのみ抽出
                            text = page.extract_text()
                            if text:
                                text_content.append(f"--- ページ {i+1} ---")
                                text_content.append(text)
                    
                    if num_pages > 10:
                        text_content.append(f"\n...(残り {num_pages - 10} ページは省略)...")
                
                return f"{file_info}\n\n" + "\n".join(text_content)
                
            except Exception as e:
                logger.error(f"PyPDF2でのPDF抽出中にエラー: {str(e)}")
                logger.error(traceback.format_exc())
                # 外部コマンドによるフォールバックを試みる
                return self._extract_pdf_fallback(file_path, file_info)
        else:
            # PyPDF2が利用できない場合はフォールバック
            return self._extract_pdf_fallback(file_path, file_info)
    
    def _extract_pdf_fallback(self, file_path, file_info):
        """PDF抽出のフォールバックメソッド（外部コマンド使用）"""
        try:
            # PowerShellを使用したPDF情報抽出（テキスト抽出なし）
            cmd = f"""
            try {{
                # ファイル情報の取得
                "PDF名: {os.path.basename(file_path)}"
                "ファイルサイズ: " + (Get-Item "{file_path}").Length + " bytes"
                "最終更新日時: " + (Get-Item "{file_path}").LastWriteTime
                "----------------------------------------"
                "このPDFからのテキスト抽出はPyPDF2ライブラリがインストールされていないため利用できません。"
                "pip install PyPDF2 でインストールしてください。"
            }} catch {{
                "エラーが発生しました: $_"
            }}
            """
            process = subprocess.Popen(
                ["powershell", "-ExecutionPolicy", "Bypass", "-Command", cmd],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True
            )
            stdout, stderr = process.communicate()
            
            if stderr:
                logger.warning(f"PDF情報取得中のエラー: {stderr}")
            
            return f"{file_info}\n\n{stdout}"
        except Exception as e:
            logger.error(f"PDF抽出フォールバック中にエラー: {str(e)}")
            return f"{file_info}\n\nPDF抽出エラー: {str(e)}"

    def _extract_docx(self, file_path):
        """Word文書(docx)からテキストを抽出"""
        file_info = self._get_file_info(file_path)
        
        # python-docxが利用可能な場合
        if self.imports['docx']:
            try:
                import docx
                
                doc = docx.Document(file_path)
                
                # 文書情報
                core_properties = doc.core_properties
                text_content = []
                
                if core_properties:
                    text_content.append(f"タイトル: {core_properties.title or '不明'}")
                    text_content.append(f"作成者: {core_properties.author or '不明'}")
                    text_content.append(f"最終更新者: {core_properties.last_modified_by or '不明'}")
                
                # 段落から本文テキストを抽出
                text_content.append("----------------------------------------")
                text_content.append("文書内容:")
                
                for i, para in enumerate(doc.paragraphs):
                    if para.text.strip():
                        text_content.append(para.text)
                
                # テーブルの内容を抽出
                if doc.tables:
                    text_content.append("\n--- テーブル内容 ---")
                    for i, table in enumerate(doc.tables):
                        if i < 5:  # 最初の5つのテーブルのみ処理
                            text_content.append(f"テーブル {i+1}:")
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    row_text.append(cell.text.strip())
                                text_content.append(" | ".join(row_text))
                            text_content.append("")
                
                return f"{file_info}\n\n" + "\n".join(text_content)
                
            except Exception as e:
                logger.error(f"python-docxでのWord抽出中にエラー: {str(e)}")
                # 外部コマンドによるフォールバックを試みる
                return self._extract_docx_fallback(file_path, file_info)
        else:
            # python-docxが利用できない場合はフォールバック
            return self._extract_docx_fallback(file_path, file_info)
    
    def _extract_docx_fallback(self, file_path, file_info):
        """Word抽出のフォールバックメソッド（外部コマンド使用）"""
        try:
            # PowerShellを使用したWord文書の情報抽出
            cmd = f"""
            try {{
                # ファイル情報の取得
                "Word文書名: {os.path.basename(file_path)}"
                "ファイルサイズ: " + (Get-Item "{file_path}").Length + " bytes"
                "最終更新日時: " + (Get-Item "{file_path}").LastWriteTime
                "----------------------------------------"
                "このWord文書からのテキスト抽出はpython-docxライブラリがインストールされていないため利用できません。"
                "pip install python-docx でインストールしてください。"
            }} catch {{
                "エラーが発生しました: $_"
            }}
            """
            process = subprocess.Popen(
                ["powershell", "-ExecutionPolicy", "Bypass", "-Command", cmd],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True
            )
            stdout, stderr = process.communicate()
            
            if stderr:
                logger.warning(f"Word情報取得中のエラー: {stderr}")
            
            return f"{file_info}\n\n{stdout}"
        except Exception as e:
            logger.error(f"Word抽出フォールバック中にエラー: {str(e)}")
            return f"{file_info}\n\nWord抽出エラー: {str(e)}"

    def _extract_xlsx(self, file_path):
        """Excelファイル(xlsx)からデータを抽出"""
        file_info = self._get_file_info(file_path)
        
        # openpyxlが利用可能な場合
        if self.imports['xlsx']:
            try:
                import openpyxl
                
                workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                
                text_content = []
                text_content.append(f"ブック名: {os.path.basename(file_path)}")
                text_content.append(f"シート数: {len(workbook.sheetnames)}")
                text_content.append(f"シート一覧: {', '.join(workbook.sheetnames)}")
                text_content.append("----------------------------------------")
                
                # 各シートの内容を抽出
                for sheet_name in workbook.sheetnames[:5]:  # 最初の5シートのみ処理
                    sheet = workbook[sheet_name]
                    text_content.append(f"--- シート: {sheet_name} ---")
                    
                    row_count = 0
                    for row in sheet.iter_rows(max_row=50):  # 最初の50行のみ処理
                        row_data = []
                        for cell in row:
                            row_data.append(str(cell.value if cell.value is not None else ""))
                        text_content.append("\t".join(row_data))
                        row_count += 1
                    
                    if row_count == 50:
                        text_content.append("...(以降省略)...")
                    
                    text_content.append("")
                
                workbook.close()
                return f"{file_info}\n\n" + "\n".join(text_content)
                
            except Exception as e:
                logger.error(f"openpyxlでのExcel抽出中にエラー: {str(e)}")
                # 外部コマンドによるフォールバックを試みる
                return self._extract_xlsx_fallback(file_path, file_info)
        else:
            # openpyxlが利用できない場合はフォールバック
            return self._extract_xlsx_fallback(file_path, file_info)
    
    def _extract_xlsx_fallback(self, file_path, file_info):
        """Excel抽出のフォールバックメソッド（外部コマンド使用）"""
        try:
            # PowerShellを使用したExcelファイルの情報抽出
            cmd = f"""
            try {{
                # ファイル情報の取得
                "Excel名: {os.path.basename(file_path)}"
                "ファイルサイズ: " + (Get-Item "{file_path}").Length + " bytes"
                "最終更新日時: " + (Get-Item "{file_path}").LastWriteTime
                "----------------------------------------"
                "このExcelからのデータ抽出はopenpyxlライブラリがインストールされていないため利用できません。"
                "pip install openpyxl でインストールしてください。"
            }} catch {{
                "エラーが発生しました: $_"
            }}
            """
            process = subprocess.Popen(
                ["powershell", "-ExecutionPolicy", "Bypass", "-Command", cmd],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True
            )
            stdout, stderr = process.communicate()
            
            if stderr:
                logger.warning(f"Excel情報取得中のエラー: {stderr}")
            
            return f"{file_info}\n\n{stdout}"
        except Exception as e:
            logger.error(f"Excel抽出フォールバック中にエラー: {str(e)}")
            return f"{file_info}\n\nExcel抽出エラー: {str(e)}"

    def _extract_pptx(self, file_path):
        """PowerPointファイル(pptx)からテキストを抽出"""
        file_info = self._get_file_info(file_path)
        
        # python-pptxが利用可能な場合
        if self.imports['pptx']:
            try:
                import pptx
                
                presentation = pptx.Presentation(file_path)
                
                text_content = []
                text_content.append(f"プレゼンテーション名: {os.path.basename(file_path)}")
                text_content.append(f"スライド数: {len(presentation.slides)}")
                text_content.append("----------------------------------------")
                
                # 各スライドからテキストを抽出
                for i, slide in enumerate(presentation.slides):
                    if i < 20:  # 最初の20スライドのみ処理
                        text_content.append(f"--- スライド {i+1} ---")
                        
                        # スライドのタイトル
                        if slide.shapes.title:
                            text_content.append(f"タイトル: {slide.shapes.title.text}")
                        
                        # スライド内のテキスト要素を抽出
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    text_content.append(paragraph.text)
                        
                        text_content.append("")
                
                if len(presentation.slides) > 20:
                    text_content.append(f"...(残り {len(presentation.slides) - 20} スライドは省略)...")
                
                return f"{file_info}\n\n" + "\n".join(text_content)
                
            except Exception as e:
                logger.error(f"python-pptxでのPowerPoint抽出中にエラー: {str(e)}")
                # 外部コマンドによるフォールバックを試みる
                return self._extract_pptx_fallback(file_path, file_info)
        else:
            # python-pptxが利用できない場合はフォールバック
            return self._extract_pptx_fallback(file_path, file_info)
    
    def _extract_pptx_fallback(self, file_path, file_info):
        """PowerPoint抽出のフォールバックメソッド（外部コマンド使用）"""
        try:
            # PowerShellを使用したPowerPointファイルの情報抽出
            cmd = f"""
            try {{
                # ファイル情報の取得
                "PowerPoint名: {os.path.basename(file_path)}"
                "ファイルサイズ: " + (Get-Item "{file_path}").Length + " bytes"
                "最終更新日時: " + (Get-Item "{file_path}").LastWriteTime
                "----------------------------------------"
                "このPowerPointからのテキスト抽出はpython-pptxライブラリがインストールされていないため利用できません。"
                "pip install python-pptx でインストールしてください。"
            }} catch {{
                "エラーが発生しました: $_"
            }}
            """
            process = subprocess.Popen(
                ["powershell", "-ExecutionPolicy", "Bypass", "-Command", cmd],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True
            )
            stdout, stderr = process.communicate()
            
            if stderr:
                logger.warning(f"PowerPoint情報取得中のエラー: {stderr}")
            
            return f"{file_info}\n\n{stdout}"
        except Exception as e:
            logger.error(f"PowerPoint抽出フォールバック中にエラー: {str(e)}")
            return f"{file_info}\n\nPowerPoint抽出エラー: {str(e)}"

    def _get_file_info(self, file_path):
        """ファイルの基本情報を取得"""
        try:
            file_stats = os.stat(file_path)
            file_size = file_stats.st_size
            modified_time = datetime.fromtimestamp(file_stats.st_mtime).strftime('%Y年%m月%d日 %H:%M:%S')
            
            # ファイルサイズを適切な単位に変換
            if file_size < 1024:
                size_str = f"{file_size} bytes"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.1f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.1f} MB"
            
            return f"ファイル名: {os.path.basename(file_path)}\nファイルサイズ: {size_str}\n最終更新日時: {modified_time}"
            
        except Exception as e:
            logger.error(f"ファイル情報取得中にエラー: {str(e)}")
            return f"ファイル名: {os.path.basename(file_path)}"
